#!/usr/bin/env python3
"""
Generate a PowerPoint presentation about the tuya-convert project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    slide1_layout = prs.slide_layouts[6]  # Blank layout
    slide1 = prs.slides.add_slide(slide1_layout)

    # Add title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)
    title_box = slide1.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "TUYA-CONVERT"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    left = Inches(1)
    top = Inches(3.8)
    width = Inches(8)
    height = Inches(1)
    subtitle_box = slide1.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Liberating Smart Home Devices from the Cloud"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Purpose
    slide2_layout = prs.slide_layouts[6]  # Blank layout
    slide2 = prs.slides.add_slide(slide2_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide2.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Purpose"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Add content
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)
    content_box = slide2.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    # Add bullet points
    bullet_points = [
        "Address security vulnerabilities in Tuya's cloud-connected smart home devices",
        "Free devices from mandatory cloud dependency and proprietary firmware",
        "Enable users to flash alternative open-source firmware without soldering",
        "Provide an accessible solution for DIY smart home enthusiasts",
        "Build on security research by VTRUST, presented at 35C3 conference"
    ]

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.level = 0
        p.space_before = Pt(12)

    # Slide 3: Goals
    slide3_layout = prs.slide_layouts[6]  # Blank layout
    slide3 = prs.slides.add_slide(slide3_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide3.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Goals"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Add content
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)
    content_box = slide3.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    # Add bullet points
    bullet_points = [
        "Provide Over-The-Air (OTA) flashing for ESP8266/85-based Tuya devices",
        "Enable installation of alternative firmwares (Tasmota, ESPurna, etc.)",
        "Create automatic backups of original firmware before flashing",
        "Build a community-maintained database of compatible devices",
        "Support cross-platform deployment (Raspberry Pi, Linux, Docker)",
        "Maintain compatibility despite manufacturer firmware updates"
    ]

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.level = 0
        p.space_before = Pt(12)

    # Save presentation
    prs.save('tuya-convert-presentation.pptx')
    print("Presentation created successfully: tuya-convert-presentation.pptx")

if __name__ == "__main__":
    create_presentation()
